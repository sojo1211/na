# -*- coding: utf-8 -*-
# YOLOX 기반 손동작 인맺기 — 학술 발표 PPT
# 참고 디자인: 흰 배경 + 네이비 포인트 + 깔끔한 학술 스타일
from pptx import Presentation
from pptx.util import Pt, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from lxml import etree

# ── 색상 ────────────────────────────────────────────────────
C_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
C_NAVY   = RGBColor(0x1F, 0x38, 0x64)   # 진한 네이비 (참고 PPT 포인트색)
C_NAVY2  = RGBColor(0x2E, 0x54, 0x96)   # 중간 네이비
C_BLACK  = RGBColor(0x1A, 0x1A, 0x1A)
C_GRAY   = RGBColor(0x4D, 0x4D, 0x4D)
C_LGRAY  = RGBColor(0xF2, 0xF5, 0xF9)   # 연한 파란 회색 배경박스
C_BORDER = RGBColor(0xCC, 0xD6, 0xE8)   # 박스 테두리
C_RED    = RGBColor(0xC0, 0x00, 0x00)   # 강조 빨강
C_ORANGE = RGBColor(0xE8, 0x6A, 0x10)   # 서브 강조

prs = Presentation()
prs.slide_width  = Cm(33.87)
prs.slide_height = Cm(19.05)
W = prs.slide_width
H = prs.slide_height
BLANK = prs.slide_layouts[6]


# ── 기본 도형/텍스트 헬퍼 ────────────────────────────────────
def rect(sl, x, y, w, h, color, line_color=None):
    sp = sl.shapes.add_shape(1, x, y, w, h)
    sp.fill.solid()
    sp.fill.fore_color.rgb = color
    if line_color:
        sp.line.color.rgb = line_color
        sp.line.width = Pt(1)
    else:
        sp.line.fill.background()
    return sp


def tb(sl, text, x, y, w, h,
       size=18, bold=False, italic=False,
       color=None, align=PP_ALIGN.LEFT, wrap=True):
    t = sl.shapes.add_textbox(x, y, w, h)
    tf = t.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.name = "맑은 고딕"
    r.font.color.rgb = color or C_BLACK
    return t


def multi_tb(sl, lines, x, y, w, h, size=15, line_spacing=1.4):
    """[(text, bold, color), ...] 여러 줄 텍스트박스"""
    t = sl.shapes.add_textbox(x, y, w, h)
    tf = t.text_frame
    tf.word_wrap = True
    for i, (text, bold, color) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_before = Pt(3)
        r = p.add_run()
        r.text = text
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.name = "맑은 고딕"
        r.font.color.rgb = color or C_BLACK
    return t


# ── 슬라이드 공통 요소 ────────────────────────────────────────
def white_bg(sl):
    rect(sl, 0, 0, W, H, C_WHITE)


def title_block(sl, text, y=Cm(1.3)):
    """참고PPT 스타일: 왼쪽 네이비 세로바 + 제목 + 가로선"""
    rect(sl, Cm(1.2), y, Cm(0.25), Cm(1.15), C_NAVY)     # 세로 바
    tb(sl, text, Cm(1.7), y, W - Cm(2.5), Cm(1.15),
       size=26, bold=True, color=C_NAVY)
    rect(sl, Cm(1.2), y + Cm(1.2), W - Cm(2.4), Cm(0.06), C_NAVY)  # 가로선


def page_num(sl, n, total=10):
    tb(sl, f"{n} / {total}",
       W - Cm(3), H - Cm(0.8), Cm(2.7), Cm(0.6),
       size=11, color=C_GRAY, align=PP_ALIGN.RIGHT)


def section_label(sl, text):
    """상단 우측 섹션 배지"""
    rect(sl, W - Cm(6), Cm(1.3), Cm(4.8), Cm(0.65), C_NAVY)
    tb(sl, text, W - Cm(6), Cm(1.3), Cm(4.8), Cm(0.65),
       size=12, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)


def numbered_badge(sl, n, x, y, r=Cm(0.5)):
    """번호 다이아몬드 배지"""
    rect(sl, x, y, Cm(0.9), Cm(0.9), C_NAVY)
    tb(sl, str(n), x, y, Cm(0.9), Cm(0.9),
       size=14, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)


_HYPERLINK_REL = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships/hyperlink'
_A_NS         = 'http://schemas.openxmlformats.org/drawingml/2006/main'
_R_NS         = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'

def add_hyperlink(sl, shape, url):
    """shape 클릭 시 외부 URL을 여는 하이퍼링크를 추가한다."""
    rId = sl.part.relate_to(url, _HYPERLINK_REL, is_external=True)
    cNvPr = shape.element[0][0]   # p:sp > p:nvSpPr > p:cNvPr
    hl = etree.SubElement(cNvPr, f'{{{_A_NS}}}hlinkClick')
    hl.set(f'{{{_R_NS}}}id', rId)


# ══════════════════════════════════════════════════════════════
# SLIDE 1 — 표지
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
white_bg(sl)

# 상단 헤더 (참고PPT 스타일)
rect(sl, Cm(1.2), Cm(0.7), Cm(0.25), Cm(0.8), C_NAVY)
tb(sl, "용인대학교  AI학과  발표자료",
   Cm(1.7), Cm(0.72), W - Cm(3), Cm(0.7),
   size=13, color=C_GRAY)
rect(sl, Cm(1.2), Cm(1.55), W * 0.45, Cm(0.05), C_NAVY)

# 메인 타이틀
tb(sl, "YOLOX 기반\n손동작 인맺기",
   Cm(2.5), Cm(3.5), W - Cm(5), Cm(5.5),
   size=52, bold=True, color=C_NAVY)

# 영문 서브타이틀
tb(sl, "Hand Sign Detection with YOLOX-Nano",
   Cm(2.7), Cm(9.0), W - Cm(5), Cm(1.0),
   size=18, italic=True, color=C_GRAY)

# 구분선
rect(sl, Cm(2.5), Cm(10.4), Cm(12), Cm(0.07), C_NAVY2)

# 저자/소속 정보
tb(sl, "용인대학교",
   Cm(2.5), Cm(11.0), Cm(14), Cm(0.8),
   size=15, color=C_GRAY)
tb(sl, "오성준  ·  이예성",
   Cm(2.5), Cm(11.8), Cm(14), Cm(1.2),
   size=26, bold=True, color=C_NAVY)

# 우측 네이비 세로 강조선 (참고PPT 오른쪽 바 스타일)
rect(sl, W - Cm(1.5), Cm(10.5), Cm(0.35), Cm(5), C_NAVY)

# 우측 기술 뱃지
for i, tech in enumerate(["YOLOX-Nano", "ONNX Runtime", "OpenCV"]):
    by = Cm(11.2) + Cm(i * 1.3)
    rect(sl, W - Cm(7), by, Cm(5.2), Cm(1.0), C_LGRAY, C_BORDER)
    tb(sl, tech, W - Cm(7), by, Cm(5.2), Cm(1.0),
       size=14, bold=True, color=C_NAVY2, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════
# SLIDE 2 — 목차
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
white_bg(sl)
title_block(sl, "Contents")
section_label(sl, "목차")

contents = [
    ("연구 배경",   "YOLO vs YOLOX — 앵커 박스 개념 이해"),
    ("주제 선정",   "나루토 인맺기 아이디어 소개"),
    ("기술 스택",   "데이터셋 구성 / 모델 선택 / 하이퍼파라미터"),
    ("시스템 구현", "실시간 데모 결과 / 최적화 사례"),
]
for i, (main, sub) in enumerate(contents):
    y = Cm(3.8) + Cm(i * 3.3)
    numbered_badge(sl, i + 1, Cm(1.5), y + Cm(0.1))
    # 박스
    rect(sl, Cm(2.7), y, W - Cm(4.2), Cm(2.8), C_WHITE, C_BORDER)
    rect(sl, Cm(2.7), y, Cm(0.3), Cm(2.8), C_NAVY)  # 왼쪽 색선
    tb(sl, main, Cm(3.3), y + Cm(0.3), Cm(10), Cm(1.0),
       size=18, bold=True, color=C_NAVY)
    tb(sl, sub, Cm(3.3), y + Cm(1.3), W - Cm(5), Cm(1.1),
       size=14, color=C_GRAY)

page_num(sl, 2)

# ══════════════════════════════════════════════════════════════
# SLIDE 3 — 섹션 전환 (연구 배경)  [참고PPT "[연구배경]" 스타일]
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
rect(sl, 0, 0, W, H, C_NAVY)
# 장식선
rect(sl, 0, H // 2 - Cm(0.04), W, Cm(0.08), RGBColor(0x3A, 0x60, 0xA0))
tb(sl, "01  연구 배경",
   Cm(3), H // 2 - Cm(1.8), W - Cm(6), Cm(3),
   size=44, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
tb(sl, "Research Background",
   Cm(3), H // 2 + Cm(1.3), W - Cm(6), Cm(1.2),
   size=20, italic=True, color=RGBColor(0x99, 0xB2, 0xCC), align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════
# SLIDE 4 — 연구 배경: YOLO vs YOLOX
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
white_bg(sl)
title_block(sl, "YOLO vs YOLOX — 앵커 박스의 차이")
section_label(sl, "01  연구 배경")

# 왼쪽: 기존 YOLO
rect(sl, Cm(1.2), Cm(3.5), Cm(14.5), Cm(13.0), C_LGRAY, C_BORDER)
rect(sl, Cm(1.2), Cm(3.5), Cm(14.5), Cm(0.75), C_GRAY)
tb(sl, "기존 YOLO — 앵커 박스 방식",
   Cm(1.2), Cm(3.5), Cm(14.5), Cm(0.75),
   size=14, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

multi_tb(sl, [
    ("옷가게 비유 🧥", True, C_NAVY),
    ("S / M / L / XL / XXL — 5가지 기성복만 준비", False, C_BLACK),
    ("", False, C_BLACK),
    ("→  새 데이터 학습 시", False, C_GRAY),
    ("    앵커 크기 사전 계산 필요", False, C_GRAY),
    ("→  기존에 없는 클래스 → 설정 번거로움", False, C_RED),
], Cm(1.8), Cm(4.7), Cm(13.5), Cm(10.5), size=15)

# VS 배지
rect(sl, W // 2 - Cm(1.1), Cm(8.5), Cm(2.2), Cm(2.0), C_NAVY)
tb(sl, "VS",
   W // 2 - Cm(1.1), Cm(8.6), Cm(2.2), Cm(1.8),
   size=26, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

# 오른쪽: YOLOX
rect(sl, W // 2 + Cm(1.3), Cm(3.5), Cm(14.5), Cm(13.0), C_LGRAY, C_BORDER)
rect(sl, W // 2 + Cm(1.3), Cm(3.5), Cm(14.5), Cm(0.75), C_NAVY)
tb(sl, "YOLOX — 앵커 프리(Anchor-Free)",
   W // 2 + Cm(1.3), Cm(3.5), Cm(14.5), Cm(0.75),
   size=14, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

multi_tb(sl, [
    ("맞춤 정장 비유 👔", True, C_NAVY),
    ("손님이 오면 즉석에서 직접 측정 → 딱 맞게 제작", False, C_BLACK),
    ("", False, C_BLACK),
    ("→  사전 계산 없이 바로 학습 시작", False, C_GRAY),
    ("→  나루토 인맺기처럼 새로운 클래스도 OK", False, C_NAVY2),
    ("→  저희 프로젝트에 최적!", False, C_RED),
], W // 2 + Cm(1.9), Cm(4.7), Cm(13.5), Cm(10.5), size=15)

page_num(sl, 4)

# ══════════════════════════════════════════════════════════════
# SLIDE 5 — 섹션 전환 (주제 선정)
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
rect(sl, 0, 0, W, H, C_NAVY)
rect(sl, 0, H // 2 - Cm(0.04), W, Cm(0.08), RGBColor(0x3A, 0x60, 0xA0))
tb(sl, "02  주제 선정",
   Cm(3), H // 2 - Cm(1.8), W - Cm(6), Cm(3),
   size=44, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
tb(sl, "Topic Selection",
   Cm(3), H // 2 + Cm(1.3), W - Cm(6), Cm(1.2),
   size=20, italic=True, color=RGBColor(0x99, 0xB2, 0xCC), align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════
# SLIDE 6 — 주제 선정
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
white_bg(sl)
title_block(sl, "주제 선정")
section_label(sl, "02  주제 선정")

# 인용구 박스
rect(sl, Cm(1.2), Cm(3.5), W - Cm(2.4), Cm(2.2), C_LGRAY, C_BORDER)
rect(sl, Cm(1.2), Cm(3.5), Cm(0.3), Cm(2.2), C_ORANGE)
tb(sl, '"이 동작들을 카메라 앞에서 하면 AI가 인식할 수 있지 않을까?"',
   Cm(2.0), Cm(3.7), W - Cm(3.2), Cm(1.8),
   size=19, bold=True, italic=True, color=C_NAVY)

# 왼쪽: 선정 이유
multi_tb(sl, [
    ("▶  나루토 닌자 기술 — 인맺기(印結び)", True, C_NAVY),
    ("    누구나 어릴 때 한 번쯤 따라해본 손동작", False, C_GRAY),
    ("", False, C_BLACK),
    ("▶  주제 선정 이유", True, C_NAVY),
    ("    · 실시간 객체 탐지 & 분류 기술 직접 경험", False, C_BLACK),
    ("    · 새로운 커스텀 클래스 데이터 학습 실습", False, C_BLACK),
    ("    · 시각적으로 명확한 발표 소재", False, C_BLACK),
], Cm(1.5), Cm(6.3), W // 2 - Cm(2), Cm(10), size=16)

# 오른쪽: YouTube 영상 플레이스홀더 (클릭 → 쇼츠 재생)
VIDEO_URL = 'https://youtube.com/shorts/_BgIhGU2qF8'
vx, vy, vw, vh = W // 2 + Cm(0.5), Cm(5.8), Cm(15.2), Cm(9.6)

# 배경 — 어두운 네이비 (플레이어 느낌)
vid_bg = rect(sl, vx, vy, vw, vh, C_NAVY)
add_hyperlink(sl, vid_bg, VIDEO_URL)

# 재생 버튼 원형(사각)
btn_w = Cm(3.2)
btn_x = vx + (vw - btn_w) / 2
btn_y = vy + (vh - btn_w) / 2 - Cm(0.5)
play_btn = rect(sl, btn_x, btn_y, btn_w, btn_w, C_RED)
add_hyperlink(sl, play_btn, VIDEO_URL)
tb(sl, "▶", btn_x, btn_y, btn_w, btn_w,
   size=34, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

# 하단 레이블
tb(sl, "나루토 인맺기 쇼츠  |  YouTube Shorts",
   vx, vy + vh - Cm(1.2), vw, Cm(1.0),
   size=13, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

# 클릭 안내
tb(sl, "☞  클릭하여 영상 재생",
   vx, vy + vh + Cm(0.25), vw, Cm(0.75),
   size=12, italic=True, color=C_NAVY2, align=PP_ALIGN.CENTER)

page_num(sl, 6)

# ══════════════════════════════════════════════════════════════
# SLIDE 7 — 섹션 전환 (기술 스택)
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
rect(sl, 0, 0, W, H, C_NAVY)
rect(sl, 0, H // 2 - Cm(0.04), W, Cm(0.08), RGBColor(0x3A, 0x60, 0xA0))
tb(sl, "03  기술 스택",
   Cm(3), H // 2 - Cm(1.8), W - Cm(6), Cm(3),
   size=44, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
tb(sl, "Tech Stack",
   Cm(3), H // 2 + Cm(1.3), W - Cm(6), Cm(1.2),
   size=20, italic=True, color=RGBColor(0x99, 0xB2, 0xCC), align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════
# SLIDE 8 — 기술 스택: 데이터셋 + 모델
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
white_bg(sl)
title_block(sl, "기술 스택 — 데이터셋 & 모델")
section_label(sl, "03  기술 스택")

# 상단 통계 3개 카드
def stat_card(sl, x, val, label, sub):
    rect(sl, x, Cm(3.5), Cm(9.5), Cm(4.8), C_WHITE, C_BORDER)
    rect(sl, x, Cm(3.5), Cm(9.5), Cm(0.6), C_NAVY)
    tb(sl, label, x, Cm(3.5), Cm(9.5), Cm(0.6),
       size=12, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    tb(sl, val, x, Cm(4.3), Cm(9.5), Cm(2.2),
       size=34, bold=True, color=C_NAVY, align=PP_ALIGN.CENTER)
    tb(sl, sub, x, Cm(6.4), Cm(9.5), Cm(1.7),
       size=11, color=C_GRAY, align=PP_ALIGN.CENTER)

stat_card(sl, Cm(1.2),  "10,026장", "총 수집 이미지", "인터넷 · 직접 촬영 · 애니메이션")
stat_card(sl, Cm(12.0), "7,098장",  "라벨링 완료",   "바운딩 박스 어노테이션 완료")
stat_card(sl, Cm(22.8), "2,928장",  "라벨링 제외",   "품질 기준 미달 제외")

# 제외 기준 + 모델 선택
rect(sl, Cm(1.2), Cm(9.0), Cm(14.5), Cm(8.6), C_LGRAY, C_BORDER)
rect(sl, Cm(1.2), Cm(9.0), Cm(0.25), Cm(8.6), C_NAVY)
tb(sl, "라벨링 제외 기준", Cm(1.7), Cm(9.2), Cm(13.5), Cm(0.8),
   size=14, bold=True, color=C_NAVY)
multi_tb(sl, [
    ("① 손동작이 흐리게 촬영된 이미지", False, C_BLACK),
    ("② 손이 프레임에 절반 이상 잘린 이미지", False, C_BLACK),
    ("③ 조명이 너무 어두워 손의 윤곽이 불분명한 이미지", False, C_BLACK),
    ("", False, C_BLACK),
    ("→ 부정확한 바운딩 박스는 학습을 방해", False, C_RED),
    ("   품질 기준 통과 이미지만 학습에 사용", False, C_GRAY),
], Cm(1.7), Cm(10.2), Cm(13.5), Cm(6.5), size=14)

# YOLOX 모델 선택 표
params = [
    ("모델",        "YOLOX-Nano  (가장 가볍고 빠름)"),
    ("입력 해상도", "416 × 416"),
    ("클래스 수",   "14종"),
    ("컨피던스",    "0.7  (70% 이상 확신 시 출력)"),
    ("출력 포맷",   "ONNX  (프레임워크 독립 실행)"),
]
rect(sl, Cm(17.0), Cm(9.0), Cm(15.5), Cm(8.6), C_WHITE, C_BORDER)
rect(sl, Cm(17.0), Cm(9.0), Cm(15.5), Cm(0.6), C_NAVY)
tb(sl, "모델 & 학습 설정", Cm(17.0), Cm(9.0), Cm(15.5), Cm(0.6),
   size=12, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
for i, (k, v) in enumerate(params):
    ry = Cm(9.8) + Cm(i * 1.6)
    bg_c = C_LGRAY if i % 2 == 0 else C_WHITE
    rect(sl, Cm(17.0), ry, Cm(15.5), Cm(1.52), bg_c)
    tb(sl, k, Cm(17.2), ry + Cm(0.25), Cm(4.5), Cm(1.1),
       size=13, bold=True, color=C_NAVY)
    tb(sl, v, Cm(21.9), ry + Cm(0.25), Cm(10.5), Cm(1.1),
       size=13, color=C_BLACK)

page_num(sl, 8)

# ══════════════════════════════════════════════════════════════
# SLIDE 9 — 섹션 전환 (시스템 구현)
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
rect(sl, 0, 0, W, H, C_NAVY)
rect(sl, 0, H // 2 - Cm(0.04), W, Cm(0.08), RGBColor(0x3A, 0x60, 0xA0))
tb(sl, "04  시스템 구현",
   Cm(3), H // 2 - Cm(1.8), W - Cm(6), Cm(3),
   size=44, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
tb(sl, "System Implementation",
   Cm(3), H // 2 + Cm(1.3), W - Cm(6), Cm(1.2),
   size=20, italic=True, color=RGBColor(0x99, 0xB2, 0xCC), align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════
# SLIDE 10 — 시스템 구현: 데모 결과
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
white_bg(sl)
title_block(sl, "시스템 구현 — 실시간 데모 결과")
section_label(sl, "04  시스템 구현")

# 웹캠 프레임 모형 (흰 배경에 어울리게 네이비 테두리)
rect(sl, Cm(1.2), Cm(3.5), Cm(19.5), Cm(12.5), C_LGRAY, C_BORDER)
# 상단 바
rect(sl, Cm(1.2), Cm(3.5), Cm(19.5), Cm(0.55), C_NAVY)
tb(sl, "NARUTO HandSignDetection Simple Demo",
   Cm(1.4), Cm(3.5), Cm(19), Cm(0.55),
   size=11, color=C_WHITE)
# 내부 프레임
rect(sl, Cm(2.5), Cm(4.4), Cm(17), Cm(10.5), RGBColor(0xE8, 0xEC, 0xF5), C_BORDER)
# 바운딩 박스 모형
for dx, dy, dw, dh in [(0, 0, 14, 0.1), (0, 8.4, 14, 0.1),
                        (0, 0, 0.1, 8.5), (13.9, 0, 0.1, 8.5)]:
    rect(sl, Cm(3.5 + dx), Cm(5.0 + dy), Cm(dw), Cm(dh),
         RGBColor(0x00, 0xCC, 0x44))
tb(sl, "ID:4  U(Hare) / 卯  conf: 0.823",
   Cm(3.5), Cm(4.55), Cm(14), Cm(0.6),
   size=12, color=RGBColor(0x00, 0xAA, 0x33))
tb(sl, "Elapsed Time: 18.4ms",
   Cm(2.7), Cm(4.55), Cm(9), Cm(0.55),
   size=11, color=RGBColor(0x00, 0xAA, 0x33))
tb(sl, "[ WEBCAM DEMO ]",
   Cm(3.5), Cm(8.8), Cm(14), Cm(2),
   size=22, color=C_BORDER, align=PP_ALIGN.CENTER)

# 우측 결과 설명
multi_tb(sl, [
    ("▶  실시간 인식 결과", True, C_NAVY),
    ("", False, C_BLACK),
    ("신뢰도: 0.7 ~ 0.9 수준", False, C_BLACK),
    ("(최대 1.0 기준 — 높은 정확도)", False, C_GRAY),
    ("", False, C_BLACK),
    ("처리 속도: 약 18~25ms", False, C_BLACK),
    ("", False, C_BLACK),
    ("▶  토끼(卯) 최적화 사례", True, C_NAVY),
    ("", False, C_BLACK),
    ("다른 인맺기와 유사한 모양", False, C_BLACK),
    ("→ 기본 0.7에서 인식 누락", False, C_RED),
    ("→ 임계값 하향 → 정상 인식", False, C_NAVY2),
    ("", False, C_BLACK),
    ("임계값은 유연하게 튜닝 가능", False, C_GRAY),
], Cm(22.2), Cm(3.7), Cm(10.5), Cm(12), size=14)

page_num(sl, 10)

# ══════════════════════════════════════════════════════════════
# SLIDE 11 — 결론 & Q&A
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
white_bg(sl)
title_block(sl, "결론 & Q&A")

# 성과 카드 4개
achievements = [
    ("YOLOX-Nano\n실시간 인식 성공",    "신뢰도 0.7~0.9\n높은 정확도 달성"),
    ("앵커 프리 장점\n직접 확인",        "새 클래스 학습 시\n사전 계산 불필요"),
    ("토끼(卯) 케이스\n임계값 튜닝 해결", "엣지케이스도\n유연하게 대응"),
    ("ONNX 포맷\n환경 독립 배포",        "딥러닝 프레임워크\n없이도 실행 가능"),
]
for i, (title_t, sub_t) in enumerate(achievements):
    x = Cm(1.2) + Cm(i * 8.0)
    rect(sl, x, Cm(3.6), Cm(7.5), Cm(9.5), C_WHITE, C_BORDER)
    rect(sl, x, Cm(3.6), Cm(7.5), Cm(0.55), C_NAVY)
    rect(sl, x, Cm(3.6), Cm(0.25), Cm(9.5), C_NAVY)
    tb(sl, f"✔  {title_t}",
       x + Cm(0.5), Cm(4.4), Cm(6.5), Cm(2.8),
       size=15, bold=True, color=C_NAVY)
    tb(sl, sub_t,
       x + Cm(0.5), Cm(7.4), Cm(6.5), Cm(2.8),
       size=13, color=C_GRAY)

# 하단 감사 메시지
rect(sl, Cm(1.2), Cm(14.0), W - Cm(2.4), Cm(3.6), C_LGRAY, C_BORDER)
rect(sl, Cm(1.2), Cm(14.0), Cm(0.25), Cm(3.6), C_ORANGE)
tb(sl, "감사합니다  —  Q & A",
   Cm(2.0), Cm(14.5), W - Cm(3.5), Cm(1.5),
   size=28, bold=True, color=C_NAVY, align=PP_ALIGN.CENTER)
tb(sl, "용인대학교  오성준  ·  이예성",
   Cm(2.0), Cm(16.0), W - Cm(3.5), Cm(1.0),
   size=16, color=C_GRAY, align=PP_ALIGN.CENTER)

page_num(sl, 11, 11)

# ══════════════════════════════════════════════════════════════
OUT = r"C:\Users\sungj\OneDrive\Desktop\나루토\발표자료_나루토인맺기.pptx"
prs.save(OUT)
print("저장 완료:", OUT)
